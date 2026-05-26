"""
Agentic Pipeline Service
- Document analysis using Google Vertex AI Gemini (primary) or Anthropic Claude (fallback)
- Content extraction and structuring
- Slide generation from document content
- Q&A tutor mode (avatar answers questions within document context)

On Cloud Run, Vertex AI uses Application Default Credentials (ADC) — no API key needed.
Locally, set GEMINI_API_KEY for the google-generativeai SDK fallback.
"""
import json
import logging
import os
import re
from typing import Optional

from backend.config import ANTHROPIC_API_KEY

logger = logging.getLogger(__name__)

# ── LLM availability ──────────────────────────────────────────

GCP_PROJECT_ID = os.getenv("GCP_PROJECT_ID", "project-987f80c5-14e3-450d-9b0")
GCP_REGION = os.getenv("GCP_REGION", "asia-south1")
GEMINI_API_KEY = os.getenv("GEMINI_API_KEY", "").strip()

# Try Vertex AI first (works on Cloud Run with ADC)
HAS_VERTEXAI = False
_vertexai_model = None
try:
    import vertexai
    from vertexai.generative_models import GenerativeModel
    vertexai.init(project=GCP_PROJECT_ID, location=GCP_REGION)
    _vertexai_model = GenerativeModel("gemini-2.0-flash-001")
    HAS_VERTEXAI = True
    logger.info("Vertex AI Gemini configured via ADC")
except Exception as e:
    logger.warning(f"Vertex AI not available: {e}")

# Fallback: google-generativeai with API key (for local dev)
HAS_GENAI = False
try:
    import google.generativeai as genai
    if GEMINI_API_KEY:
        genai.configure(api_key=GEMINI_API_KEY)
        HAS_GENAI = True
        logger.info("google-generativeai configured with API key")
except ImportError:
    pass

HAS_ANTHROPIC = False
try:
    import anthropic
    if ANTHROPIC_API_KEY:
        HAS_ANTHROPIC = True
except ImportError:
    pass

HAS_LLM = HAS_VERTEXAI or HAS_GENAI or HAS_ANTHROPIC
logger.info(f"LLM availability: VertexAI={HAS_VERTEXAI}, GenAI={HAS_GENAI}, Anthropic={HAS_ANTHROPIC}")


def _llm_generate(prompt: str, max_tokens: int = 4000) -> str | None:
    """Call best available LLM. Returns response text or None."""
    # 1. Vertex AI (Cloud Run ADC)
    if HAS_VERTEXAI:
        try:
            response = _vertexai_model.generate_content(prompt)
            return response.text.strip()
        except Exception as e:
            logger.error(f"Vertex AI failed: {e}")

    # 2. google-generativeai (API key, local dev)
    if HAS_GENAI:
        try:
            model = genai.GenerativeModel("gemini-2.0-flash")
            response = model.generate_content(prompt)
            return response.text.strip()
        except Exception as e:
            logger.error(f"google-generativeai failed: {e}")

    # 3. Anthropic Claude
    if HAS_ANTHROPIC:
        try:
            client = anthropic.Anthropic(api_key=ANTHROPIC_API_KEY)
            message = client.messages.create(
                model="claude-sonnet-4-20250514",
                max_tokens=max_tokens,
                messages=[{"role": "user", "content": prompt}],
            )
            return message.content[0].text.strip()
        except Exception as e:
            logger.error(f"Claude failed: {e}")

    return None


# ── Document parsers ──────────────────────────────────────────

try:
    from PyPDF2 import PdfReader
    HAS_PDF = True
except ImportError:
    HAS_PDF = False

try:
    from pptx import Presentation as PptxPresentation
    HAS_PPTX = True
except ImportError:
    HAS_PPTX = False

try:
    from docx import Document as DocxDocument
    HAS_DOCX = True
except ImportError:
    HAS_DOCX = False


# ── Text extraction ───────────────────────────────────────────

def extract_text_from_file(file_path: str) -> str:
    """Extract text content from uploaded document."""
    ext = os.path.splitext(file_path)[1].lower()
    logger.info(f"Extracting text from {ext} file: {os.path.basename(file_path)}")

    text = ""
    if ext == '.pdf' and HAS_PDF:
        text = _extract_pdf(file_path)
    elif ext == '.pptx' and HAS_PPTX:
        text = _extract_pptx(file_path)
    elif ext == '.docx' and HAS_DOCX:
        text = _extract_docx(file_path)
    else:
        logger.warning(f"No parser for {ext} (PDF={HAS_PDF}, PPTX={HAS_PPTX}, DOCX={HAS_DOCX})")

    logger.info(f"Extracted {len(text)} characters")
    return text


def _extract_pdf(file_path: str) -> str:
    text_parts = []
    reader = PdfReader(file_path)
    for page in reader.pages[:50]:
        page_text = page.extract_text()
        if page_text:
            text_parts.append(page_text.strip())
    return "\n\n".join(text_parts)


def _extract_pptx(file_path: str) -> str:
    text_parts = []
    prs = PptxPresentation(file_path)
    for slide_num, slide in enumerate(prs.slides[:50], 1):
        slide_text = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                for paragraph in shape.text_frame.paragraphs:
                    text = paragraph.text.strip()
                    if text:
                        slide_text.append(text)
        if slide_text:
            text_parts.append(f"[Slide {slide_num}]\n" + "\n".join(slide_text))
    return "\n\n".join(text_parts)


def _extract_docx(file_path: str) -> str:
    text_parts = []
    doc = DocxDocument(file_path)
    for para in doc.paragraphs[:200]:
        text = para.text.strip()
        if text:
            text_parts.append(text)
    return "\n\n".join(text_parts)


# ── Slide generation ─────────────────────────────────────────

SLIDE_PROMPT = """You are a professional AI teaching avatar creating a presentation from a document.
Your goal is to produce a presentation that feels like a real human educator explaining this topic in a conversational, engaging way — NOT a robotic recitation.

Analyze the following document and create a professional presentation with 8-12 slides.
Each slide should have a title and content (either a paragraph or bullet points).

IMPORTANT GUIDELINES:
- Write content as if a friendly, knowledgeable teacher is speaking to students
- Use conversational language: "Let's look at...", "Here's something interesting...", "Notice how..."
- Make bullet points feel like talking points, not dry facts
- Include transitions between ideas
- The tone should be warm, clear, and educational
- Do NOT generate sexual, violent, hateful, or harmful content under any circumstances
- Stay strictly within the document's content — do not invent information

Return ONLY a JSON array of slides in this exact format:
[
  {{"title": "Slide Title", "body": "Optional paragraph text or null", "bullets": ["bullet 1", "bullet 2"] or null, "type": "title|content|summary"}},
  ...
]

Rules:
- First slide should be a title slide with type "title"
- Last slide should be a summary/conclusion with type "summary"
- Middle slides should be type "content"
- Each slide should have EITHER body OR bullets (not both)
- Bullet points should be concise (under 15 words each)
- Keep content clear, professional, and engaging
- Extract key insights, facts, and takeaways from the document
- DO NOT make up information — only use content from the document

Document Title: {title}

Document Content:
{doc_excerpt}"""


def generate_slides_from_content(document_text: str, title: str) -> list:
    """Generate structured presentation slides from document content."""
    logger.info(f"generate_slides: text_len={len(document_text)}, HAS_LLM={HAS_LLM}")

    if HAS_LLM and document_text.strip():
        max_chars = 15000
        doc_excerpt = document_text[:max_chars]
        if len(document_text) > max_chars:
            doc_excerpt += "\n\n[... document continues ...]"

        prompt = SLIDE_PROMPT.format(title=title, doc_excerpt=doc_excerpt)
        response = _llm_generate(prompt)

        if response:
            logger.info(f"LLM response length: {len(response)} chars")
            slides = _parse_slides_json(response)
            if slides:
                logger.info(f"Generated {len(slides)} slides from LLM")
                return slides
            else:
                logger.warning("Failed to parse slides from LLM response")

    logger.warning("Falling back to heuristic slide generation")
    return _generate_heuristic(document_text, title)


def _parse_slides_json(response_text: str) -> list | None:
    """Parse JSON slides from an LLM response. Returns None on failure."""
    json_match = re.search(r'\[[\s\S]*\]', response_text)
    if not json_match:
        return None
    try:
        slides = json.loads(json_match.group())
        valid_slides = []
        for s in slides:
            if isinstance(s, dict) and "title" in s:
                valid_slides.append({
                    "title": str(s.get("title", "")),
                    "body": s.get("body"),
                    "bullets": s.get("bullets") if isinstance(s.get("bullets"), list) else None,
                    "type": s.get("type", "content"),
                })
        return valid_slides if valid_slides else None
    except (json.JSONDecodeError, ValueError) as e:
        logger.error(f"Failed to parse slides JSON: {e}")
        return None


def _generate_heuristic(document_text: str, title: str) -> list:
    """Generate slides using text heuristics when APIs are unavailable."""
    slides = [
        {"title": title, "body": "AI-Generated Presentation", "bullets": None, "type": "title"}
    ]

    if not document_text.strip():
        slides.extend([
            {"title": "Overview", "body": "This presentation covers the key topics from your uploaded document.", "bullets": None, "type": "content"},
            {"title": "Key Points", "body": None, "bullets": [
                "Content analysis in progress",
                "AI-powered slide generation",
                "Professional avatar delivery",
                "Interactive Q&A available",
            ], "type": "content"},
            {"title": "Summary", "body": "Your AI presenter is ready to deliver this content to your team.", "bullets": None, "type": "summary"},
        ])
        return slides

    paragraphs = [p.strip() for p in document_text.split('\n\n') if p.strip() and len(p.strip()) > 20]
    chunk_size = max(1, len(paragraphs) // 8)
    chunks = []
    for i in range(0, min(len(paragraphs), 40), chunk_size):
        chunk = paragraphs[i:i + chunk_size]
        chunks.append("\n".join(chunk))
        if len(chunks) >= 8:
            break

    for idx, chunk in enumerate(chunks):
        lines = chunk.split('\n')
        slide_title = lines[0][:60] if lines else f"Section {idx + 1}"
        if len(slide_title) > 50:
            slide_title = slide_title[:47] + "..."
        remaining_lines = lines[1:] if len(lines) > 1 else lines
        if len(remaining_lines) >= 3:
            bullets = [line[:80] for line in remaining_lines[:6]]
            slides.append({"title": slide_title, "body": None, "bullets": bullets, "type": "content"})
        else:
            body = " ".join(remaining_lines)[:250]
            slides.append({"title": slide_title, "body": body or None, "bullets": None, "type": "content"})

    slides.append({
        "title": "Summary & Next Steps",
        "body": None,
        "bullets": [
            "Key insights have been extracted from your document",
            "Ask your AI tutor any questions about this content",
            "Share this presentation with your team",
            "Track engagement through analytics",
        ],
        "type": "summary",
    })
    return slides


# ── Q&A tutor mode ────────────────────────────────────────────

QA_SYSTEM_PROMPT = """You are an AI teaching avatar presenting a document titled "{title}".
You speak like a warm, knowledgeable human educator — conversational, clear, and engaging.
This should feel like a live discussion, NOT a robotic Q&A session.

Your role:
1. Answer questions about the document content accurately and naturally
2. Explain concepts as if you're having a friendly conversation with a student
3. Use phrases like "Great question!", "Let me explain that...", "Think of it this way..."
4. Add brief analogies or examples when helpful
5. Stay within the context of the document — if asked something unrelated, gently redirect
6. NEVER generate sexual, violent, hateful, or harmful content

Keep answers concise (2-4 sentences for simple questions, up to a paragraph for complex ones).
Be the kind of teacher everyone wishes they had — patient, enthusiastic, and clear.

Document Content:
{doc_excerpt}"""


def answer_question(document_text: str, question: str, presentation_title: str, chat_history: Optional[list] = None) -> str:
    """Answer a question about the document content."""
    if HAS_LLM and document_text.strip():
        max_chars = 12000
        doc_excerpt = document_text[:max_chars]
        system = QA_SYSTEM_PROMPT.format(title=presentation_title, doc_excerpt=doc_excerpt)

        parts = [system + "\n\n"]
        if chat_history:
            for msg in chat_history[-6:]:
                role = "User" if msg["role"] == "user" else "Assistant"
                parts.append(f"{role}: {msg['content']}\n")
        parts.append(f"User: {question}\nAssistant:")

        result = _llm_generate("".join(parts))
        if result:
            return result

    return _answer_fallback(question, presentation_title)


def _answer_fallback(question: str, title: str) -> str:
    return (
        f"Thank you for your question about \"{title}\". "
        f"I'd be happy to help explain this content. The document covers several important topics. "
        f"For a detailed answer to \"{question[:50]}...\", I recommend reviewing the relevant slides "
        f"in the presentation. Is there a specific section you'd like me to focus on?"
    )
